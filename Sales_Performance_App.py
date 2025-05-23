import streamlit as st
import pandas as pd
from PIL import Image
import plotly.express as px
from prophet import Prophet
from prophet.plot import plot_plotly
from pptx import Presentation
from pptx.util import Inches
import json

# --- Sayfa Yapılandırma ---
st.set_page_config(
    page_title="Satış Analizi Dashboard",
    page_icon="Logo.jpg",
    layout="wide"
)

# --- Sidebar: Logo ve Veri Yükleme ---
st.sidebar.image("Logo.jpg", width=150)
st.sidebar.title("📂 Veri Yükle")
excel_file = st.sidebar.file_uploader("Excel dosyasını yükleyin", type=["xlsx","xls"])

if excel_file:
    # Veri Yükleme
    df_sales = pd.read_excel(excel_file, sheet_name="SATIŞ")
    df_cross = pd.read_excel(excel_file, sheet_name="ÇAPRAZ SATIŞ")
    df_demo = pd.read_excel(excel_file, sheet_name="ILCE DEMOGRAFI")

    # Tarih formatı
    df_sales['YEARMONTH'] = pd.to_datetime(df_sales['YEARMONTH'].astype(str), format='%Y%m')
    df_cross['AY'] = pd.to_datetime(df_cross['AY'].astype(str), format='%Y%m')

    # Başlık
    st.title("Satış Analizi Dashboard")
    st.markdown("Bu uygulama 2021-2022 dönemine ait satış, çapraz satış ve demografi verilerini analiz eder.")

    # Sekmeler
    tab1, tab2, tab3 = st.tabs(["Veri Önizleme","EDA & Görselleştirme","Tahmin & Rapor"])

    # -------- VERİ ÖNİZLEME --------
    with tab1:
        st.subheader("1. Ürün1 Satış Verisi")
        st.dataframe(df_sales.head(10))
        st.write(df_sales.describe())

        st.subheader("2. Ürün2 Çapraz Satış Verisi")
        st.dataframe(df_cross.head(10))
        st.write(df_cross.describe())

        st.subheader("3. Demografi Verisi")
        st.dataframe(df_demo.head(10))
        st.write(df_demo.describe())

    # -------- EDA & GÖRSELLEŞTİRME --------
    with tab2:
        st.subheader("Aylık Toplam Satış Adedi")
        df_time = df_sales.groupby('YEARMONTH')[['URUNADET','URUNHACIM']].sum().reset_index()
        fig_time = px.line(df_time, x='YEARMONTH', y=['URUNADET'], title='Zaman Serisi Analizi - Adet')
        st.plotly_chart(fig_time, use_container_width=True)

        st.subheader("Aylık Toplam Satış Hacmi")
        fig_time2 = px.line(df_time, x='YEARMONTH', y=['URUNHACIM'], title='Zaman Serisi Analizi - Hacim')
        st.plotly_chart(fig_time2, use_container_width=True)


        
        # İl Bazlı Harita (Scatter Mapbox)
        st.subheader("Şube Performansı Haritası (İl)")
        # Dealer -> City ilişkilendirme
        dealer_city = df_cross[['DEALER_CODE', 'CITY']].drop_duplicates()
        dealer_city = dealer_city.dropna()
        dealer_city['DEALER_CODE'] = dealer_city['DEALER_CODE'] // 10
        df_sales_map = pd.merge(df_sales,dealer_city, on='DEALER_CODE', how='inner')
        df_city_sales = df_sales_map.groupby('CITY')['URUNADET'].sum().reset_index()

        # GeoJSON yükleme
        try:
            with open('tr-cities.json', 'r', encoding='utf-8') as f:
                geojson_data = json.load(f)
        except FileNotFoundError:
            st.error("GeoJSON dosyası bulunamadı. 'tr-cities.json' ekleyin.")
            st.stop()

        # Her ilin centroid koordinatını hesapla
        centroids = {}
        for feat in geojson_data['features']:
            name = feat['properties']['name']
            coords = []
            geom = feat['geometry']
            if geom['type'] == 'Polygon':
                rings = geom['coordinates']
            else:
                rings = [ring for poly in geom['coordinates'] for ring in poly]
            for ring in rings:
                coords.extend(ring)
            lons = [c[0] for c in coords]
            lats = [c[1] for c in coords]
            if lats and lons:
                centroids[name] = {'lat': sum(lats)/len(lats), 'lon': sum(lons)/len(lons)}

        # Koordinatları data frame'e ekle
        df_city_sales['lat'] = df_city_sales['CITY'].map(lambda x: centroids.get(x, {}).get('lat', 38))
        df_city_sales['lon'] = df_city_sales['CITY'].map(lambda x: centroids.get(x, {}).get('lon', 35))

        # Scatter Mapbox
        fig_map = px.scatter_mapbox(
            df_city_sales,
            lat='lat',
            lon='lon',
            size='URUNADET',
            color='URUNADET',
            hover_name='CITY',
            hover_data={'URUNADET': True},
            size_max=40,
            zoom=5,
            mapbox_style='open-street-map',
            title='İl Bazında Toplam Ürün1 Satış Adedi'
        )
        st.plotly_chart(fig_map, use_container_width=True)




    # -------- TAHMİN & RAPOR --------
    with tab3:
        st.subheader("Neden Prophet Kullanıyoruz?")
        st.markdown(
            "Prophet, satış verilerindeki mevsimsellik, tatil etkileri ve regresör ekleme imkanları sayesinde öngörü doğruluğunu artırmak için seçildi."
        )
        product = st.selectbox("Tahminlenecek Ürün", ["Ürün1", "Ürün2"])
        if st.button("Tahmini Hesapla"):
            # Veri hazırlığı
            if product == "Ürün1":
                df_main = (
                    df_sales.groupby('YEARMONTH')[['URUNADET','URUNHACIM','ABONE_YAS_0_3AY','ABONE_YAS_4_12AY','ABONE_YAS_1_3YAS','ABONE_YAS_3_YAS']]
                    .sum().reset_index()
                )
                df_main = df_main.rename(columns={'YEARMONTH':'ds','URUNADET':'y'})
                regressors = ['URUNHACIM','ABONE_YAS_0_3AY','ABONE_YAS_4_12AY','ABONE_YAS_1_3YAS','ABONE_YAS_3_YAS']
            else:
                df_main = (
                    df_cross.groupby('AY')[['ÇAPRAZ ÜRÜN ADET','5GUNIPTAL','6-45GUNIPTAL']]
                    .sum().reset_index()
                )
                df_main = df_main.rename(columns={'AY':'ds','ÇAPRAZ ÜRÜN ADET':'y'})
                regressors = ['5GUNIPTAL','6-45GUNIPTAL']
            df_main['ds'] = pd.to_datetime(df_main['ds'])
            # Demografi regresörler
            demo_mean = df_demo.groupby('IL').mean(numeric_only=True).mean().to_dict()
            for k,v in demo_mean.items():
                df_main[k] = v
                regressors.append(k)

            # Prophet modeli
            m = Prophet()
            for r in regressors:
                m.add_regressor(r)
            m.fit(df_main[['ds','y']+regressors])

            # Gelecek ay
            future = m.make_future_dataframe(periods=1, freq='M')
            last = df_main.iloc[-1]
            for r in regressors:
                future[r] = last[r]
            fg = m.predict(future)

            # Görsel
            st.plotly_chart(plot_plotly(m, fg), use_container_width=True)
            pred = fg.loc[fg['ds']==fg['ds'].max(),'yhat'].iloc[0]
            st.write(f"2023 Ocak öngörü ({product}): {pred:.0f}")
            st.session_state['forecast_df'] = fg

        # ------------ Rapor İndir ------------
        st.markdown("---")
        st.subheader("Rapor İndir")
        if st.button("PowerPoint Oluştur ve İndir"):
            if 'forecast_df' not in st.session_state:
                st.error("Önce 'Tahmini Hesapla' butonuna tıklayın.")
            else:
                prs = Presentation()
                # Slide 1: Ürün1 Satış
                slide1 = prs.slides.add_slide(prs.slide_layouts[5])
                slide1.shapes.title.text = "Ürün1 Satış Verisi"
                tbl1 = df_sales.head(3)
                rows1, cols1 = tbl1.shape
                table1 = slide1.shapes.add_table(rows1+1, cols1, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
                for i, col in enumerate(tbl1.columns): table1.cell(0,i).text = str(col)
                for r_idx, (_, row) in enumerate(tbl1.iterrows(), start=1):
                    for c_idx, val in enumerate(row): table1.cell(r_idx,c_idx).text = str(val)

                # Slide 2: Ürün2 Çapraz Satış
                slide2 = prs.slides.add_slide(prs.slide_layouts[5])
                slide2.shapes.title.text = "Ürün2 Çapraz Satış Verisi"
                tbl2 = df_cross.head(3)
                rows2, cols2 = tbl2.shape
                table2 = slide2.shapes.add_table(rows2+1, cols2, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
                for i, col in enumerate(tbl2.columns): table2.cell(0,i).text = str(col)
                for r_idx, (_, row) in enumerate(tbl2.iterrows(), start=1):
                    for c_idx, val in enumerate(row): table2.cell(r_idx,c_idx).text = str(val)

                # Slide 3: Demografi (Sadece ilk 7 kolon)
                slide3 = prs.slides.add_slide(prs.slide_layouts[5])
                slide3.shapes.title.text = "Demografi Verisi"
                # Sadece ilk 7 kolona indirgenmiş tablo
                cols_to_show = df_demo.columns[:7]
                tbl3 = df_demo[cols_to_show].head(3)
                rows3, cols3 = tbl3.shape
                table3 = slide3.shapes.add_table(rows3+1, cols3, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
                for i, col in enumerate(cols_to_show):
                    table3.cell(0, i).text = str(col)
                for r_idx, (_, row) in enumerate(tbl3.iterrows(), start=1):
                    for c_idx, col in enumerate(cols_to_show):
                        table3.cell(r_idx, c_idx).text = str(row[col])


                # Slide 4: Tahmin
                slide4 = prs.slides.add_slide(prs.slide_layouts[5])
                slide4.shapes.title.text = "2023 Ocak Öngörü"
                forecast = st.session_state['forecast_df']
                df_res = forecast[['ds','yhat']].tail(1)
                rows4, cols4 = df_res.shape
                table4 = slide4.shapes.add_table(rows4+1, cols4, Inches(0.5), Inches(1.5), Inches(6), Inches(2)).table
                for i, col in enumerate(df_res.columns): table4.cell(0,i).text = str(col)
                for r_idx, (_, row) in enumerate(df_res.iterrows(), start=1):
                    for c_idx, val in enumerate(row): table4.cell(r_idx,c_idx).text = str(val)

                # Kaydet ve indir
                prs.save("sales_report.pptx")
                with open("sales_report.pptx","rb") as f:
                    st.download_button("PowerPoint İndir", f, file_name="sales_analysis.pptx")
